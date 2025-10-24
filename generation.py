import tkinter as tk
from tkinter import messagebox
import os
from docxtpl import DocxTemplate
from docx2pdf import convert
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
import openpyxl
import pandas as pd
from datetime import datetime

# Function for generating certificates

def generate_certificates(excel_file, template_file, output_folder, font_path="ariali.ttf", font_size=110):
    data = pd.read_excel(excel_file)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    font_name = ImageFont.truetype(font_path, font_size)
    for index, row in data.iterrows():
        name = row["Name"]
        certificate = Image.open(template_file)
        draw = ImageDraw.Draw(certificate)
        if len(name) >= 15 and len(name) < 25:
            name_position = (500, 600)
        elif len(name) >= 10 and len(name) < 15:
            name_position = (700, 600)
        else:
            name_position = (730, 600)
        draw.text(name_position, name, fill="gold", font=font_name)
        output_path = os.path.join(output_folder, "certificate_" + name + ".png")
        certificate.save(output_path)
        print(f"Certificate generated for {name} and saved to {output_path}")
    print("All certificates have been generated!")

    data = pd.read_excel(excel_file)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    font_name = ImageFont.truetype(font_path, font_size)
    prs = Presentation()

    for index, row in data.iterrows():
        name = row["Name"]
        certificate = Image.open(template_file)
        draw = ImageDraw.Draw(certificate)
        name_position = (500, 600) if len(name) >= 15 else (700, 600)
        draw.text(name_position, name, fill="gold", font=font_name)
        image_path = os.path.join(output_folder, f"certificate_{name}.png")
        certificate.save(image_path)
        print(f"Certificate generated for {name} and saved to {image_path}")
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8.5), height=Inches(6))

    pptx_path = os.path.join(output_folder, "certificates_presentation.pptx")
    prs.save(pptx_path)
    print(f"PowerPoint presentation saved to {pptx_path}")

    print("All certificates have been generated and added to the PowerPoint presentation!")

# Functions for generating Associate Degree documents

def AssociateExcel_data(filename):
    workbook = openpyxl.load_workbook(filename)
    sheet = workbook.active
    return list(sheet.values)

def AssociateDocument(template, output_directory, student):
    doc = DocxTemplate(template)
    current_date = datetime.now().strftime("%B %d, %Y")
    doc.render({
        'name_kh': student[2],
        'g1': student[4],
        'id_kh': student[0],
        'name_e': student[3],
        'g2': student[5],
        'id_e': student[1],
        'dob_kh': student[6],
        'pro_kh': student[8],
        'dob_e': student[7],
        'pro_e': student[9],
        'ed_kh': student[10],
        'ed_e': student[11],
        'cur_date': current_date
    })
    doc_name = os.path.join(output_directory, "{}.docx".format(student[3]))
    doc.save(doc_name)
    return doc_name

def AssociateConvertPDF(doc_path, pdf_directory):
    pdf_path = os.path.join(pdf_directory, os.path.splitext(os.path.basename(doc_path))[0] + ".pdf")
    convert(doc_path, pdf_path)
    return pdf_path

def GeneratAssociate(option):
    excel_file = "Associate.xlsx"
    template_file = "WEP Temporary Certificate - Template.docx"
    docx_directory = "Associate_Documents"
    pdf_directory = "Associate_PDF"

    os.makedirs(docx_directory, exist_ok=True)
    os.makedirs(pdf_directory, exist_ok=True)
    data_rows = AssociateExcel_data(excel_file)

    for row in data_rows[1:]:
        if option in ["doc", "both"]:
            doc_path = AssociateDocument(template_file, docx_directory, row)
        if option in ["pdf", "both"]:
            if option == "pdf":
                doc_path = AssociateDocument(template_file, pdf_directory, row)
            AssociateConvertPDF(doc_path, pdf_directory)
            if option == "pdf":
                os.remove(doc_path)
    print("All files for option '{}' have been generated!".format(option))

# Functions for generating Transcripts

def TranscriptExcel_data(filename):
    workbook = openpyxl.load_workbook(filename)
    sheet = workbook.active
    return list(sheet.values)

def TranscriptDocument(template, output_directory, row_data):
    doc = DocxTemplate(template)
    current_date = datetime.now().strftime("%B %d, %Y")
    doc.render({
        "student_id": row_data[0],
        "first_name": row_data[1],
        "last_name": row_data[2],
        "logic": row_data[3],
        "l_g": row_data[4],
        "bcum": row_data[5],
        "bc_g": row_data[6],
        "design": row_data[7],
        "d_g": row_data[8],
        "p1": row_data[9],
        "p1_g": row_data[10],
        "e1": row_data[11],
        "e1_g": row_data[12],
        "wd": row_data[13],
        "wd_g": row_data[14],
        "algo": row_data[15],
        "al_g": row_data[16],
        "p2": row_data[17],
        "p2_g": row_data[18],
        "e2": row_data[19],
        "e2_g": row_data[20],
        "sd": row_data[21],
        "sd_g": row_data[22],
        "js": row_data[23],
        "js_g": row_data[24],
        "php": row_data[25],
        "ph_g": row_data[26],
        "db": row_data[27],
        "db_g": row_data[28],
        "vc1": row_data[29],
        "v1_g": row_data[30],
        "node": row_data[31],
        "no_g": row_data[32],
        "e3": row_data[33],
        "e3_g": row_data[34],
        "p3": row_data[35],
        "p3_g": row_data[36],
        "oop": row_data[37],
        "op_g": row_data[38],
        "lar": row_data[39],
        "lar_g": row_data[40],
        "vue": row_data[41],
        "vu_g": row_data[42],
        "vc2": row_data[43],
        "v2_g": row_data[44],
        "e4": row_data[45],
        "e4_g": row_data[46],
        "p4": row_data[47],
        "p4_g": row_data[48],
        "int": row_data[49],
        "in_g": row_data[50],
        'cur_date': current_date
    })
    doc_name = os.path.join(output_directory, "{}.docx".format(row_data[1]))
    doc.save(doc_name)
    return doc_name

def TranscriptPdf(doc_path, pdf_directory):
    pdf_path = os.path.join(pdf_directory, os.path.splitext(os.path.basename(doc_path))[0] + ".pdf")
    convert(doc_path, pdf_path)
    return pdf_path

def generate_transcripts(option):
    excel_file = "data.xlsx"
    template_file = "template-pnc.docx"
    docx_directory = "Transcript_Doc"
    pdf_directory = "Transcript_PDF"

    os.makedirs(docx_directory, exist_ok=True)
    os.makedirs(pdf_directory, exist_ok=True)
    data_rows = TranscriptExcel_data(excel_file)

    for row in data_rows[1:]:
        if option in ["doc", "both"]:
            doc_path = TranscriptDocument(template_file, docx_directory, row)
        if option in ["pdf", "both"]:
            if option == "pdf":
                doc_path = TranscriptDocument(template_file, pdf_directory, row)
            TranscriptPdf(doc_path, pdf_directory)
            if option == "pdf":
                os.remove(doc_path)
    print("All files for option '{}' have been generated!".format(option))

# GUI Implementation

def create_ui():
    def show_option_menu(title, generate_function):
        option_window = tk.Toplevel(window)
        option_window.title(title)
        tk.Label(option_window, text="Select the output format:").pack(padx=20, pady=10)

        def generate_with_option(selected_option):
            generate_function(selected_option)
            result_label.config(text="{} generated successfully!".format(title))  # Update the message
            option_window.destroy()

        tk.Button(option_window, text="DOCX Only", bg="blue", fg="white", command=lambda: generate_with_option("doc")).pack(padx=20, pady=5)
        tk.Button(option_window, text="PDF Only", bg="blue", fg="white", command=lambda: generate_with_option("pdf")).pack(padx=20, pady=5)
        tk.Button(option_window, text="Both DOCX and PDF", bg="blue", fg="white", command=lambda: generate_with_option("both")).pack(padx=20, pady=5)

    def generate_certificates_direct():
        generate_certificates(
            excel_file="Certificate.xlsx",
            template_file="certificate.png",
            output_folder="Certificates"
        )
        result_label.config(text="Certificates generated successfully!")  # Update the message

    def generate_all():
        option_window = tk.Toplevel(window)
        option_window.title("Generate All")
        tk.Label(option_window, text="Select the output format for all documents:").pack(padx=20, pady=10)

        def generate_all_with_option(selected_option):
            generate_transcripts(selected_option)
            generate_certificates(
                excel_file="Certificate.xlsx",
                template_file="certificate.png",
                output_folder="Certificates"
            )
            GeneratAssociate(selected_option)
            option_window.destroy()
            result_label.config(text="All documents have been generated successfully!")  # Update the message

        tk.Button(option_window, text="DOCX Only", bg="blue", fg="white", command=lambda: generate_all_with_option("doc")).pack(padx=20, pady=5)
        tk.Button(option_window, text="PDF Only", bg="blue", fg="white", command=lambda: generate_all_with_option("pdf")).pack(padx=20, pady=5)
        tk.Button(option_window, text="Both DOCX and PDF", bg="blue", fg="white", command=lambda: generate_all_with_option("both")).pack(padx=20, pady=5)

    window = tk.Tk()
    window.geometry("500x500")
    window.title("Automated Document Generator")

    # Add a styled title

    title_label = tk.Label(
        window,
        text="Automated Document Generation",
        font=("Helvetica", 34, "bold"),
        fg="gray"
    )
    title_label.pack(pady=20)

    # Add result feedback label

    result_label = tk.Label(window, text="", fg="green", font=("Arial", 12))
    result_label.pack(pady=10)

    # Buttons

    tk.Button(window, text="Generate Transcript", bg="green", fg="white",font=("Arial", 12, "bold"), width=25, height=2, padx=18, command=lambda: show_option_menu("Transcript", generate_transcripts)).pack(pady=10)
    tk.Button(window, text="Generate Certificates", bg="green", fg="white", font=("Arial", 12, "bold"), width=25, height=2, padx=18, command=generate_certificates_direct).pack(pady=10)
    tk.Button(window, text="Generate Associate", bg="green", fg="white", font=("Arial", 12, "bold"), width=25, height=2, padx=18, command=lambda: show_option_menu("Associate Documents", GeneratAssociate)).pack(pady=10)
    tk.Button(window, text="Generate All", bg="green", fg="white", font=("Arial", 12, "bold"), width=25, height=2, padx=18, command=generate_all).pack(pady=20)

    window.mainloop()

create_ui()